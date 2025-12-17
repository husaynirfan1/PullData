"""PowerPoint formatter for PullData output.

Creates PowerPoint presentations with template-based generation,
automatic slide layouts, and table support.
"""

from io import BytesIO
from typing import Any, Dict, List, Optional

from pulldata.synthesis.base import FormatConfig, FormatterError, OutputData, OutputFormatter

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PowerPointFormatter(OutputFormatter):
    """PowerPoint formatter using python-pptx.
    
    Features:
    - Template-based generation
    - Automatic slide layouts (title, content, table, chart)
    - Table insertion
    - Text box and image support
    
    Example:
        >>> formatter = PowerPointFormatter()
        >>> data = OutputData(title="Report", content="Summary", tables=[...])
        >>> formatter.save(data, "output.pptx")
    """

    def __init__(
        self,
        config: Optional[FormatConfig] = None,
        slide_width: Optional[Inches] = None,
        slide_height: Optional[Inches] = None,
    ):
        """Initialize PowerPoint formatter.
        
        Args:
            config: Formatter configuration
            slide_width: Custom slide width (default: standard 10 inches)
            slide_height: Custom slide height (default: standard 7.5 inches)
        """
        super().__init__(config)

        if not PPTX_AVAILABLE:
            raise FormatterError(
                "python-pptx not available",
                formatter="PowerPoint",
                details={"install": "pip install python-pptx"},
            )

        self.slide_width = slide_width
        self.slide_height = slide_height

    @property
    def file_extension(self) -> str:
        return ".pptx"

    @property
    def format_name(self) -> str:
        return "PowerPoint"

    def format(self, data: OutputData) -> bytes:
        """Format data as PowerPoint presentation.
        
        Args:
            data: Data to format
            
        Returns:
            PowerPoint file as bytes
        """
        if not PPTX_AVAILABLE:
            raise FormatterError(
                "python-pptx not available",
                formatter="PowerPoint",
                details={"install": "pip install python-pptx"},
            )

        # Create presentation (use template if specified)
        if self.config.template_path:
            prs = Presentation(str(self.config.template_path))
        else:
            prs = Presentation()

        # Set custom dimensions if specified
        if self.slide_width:
            prs.slide_width = int(self.slide_width)
        if self.slide_height:
            prs.slide_height = int(self.slide_height)

        # Title slide
        self._create_title_slide(prs, data)

        # Content slide
        if data.content:
            self._create_content_slide(prs, "Summary", data.content)

        # Section slides
        if data.sections:
            for section in data.sections:
                title = section.get("title", "Section")
                content = section.get("content", "")
                self._create_content_slide(prs, title, content)

        # Table slides
        if data.tables:
            for i, table_data in enumerate(data.tables):
                table_name = table_data.get("name", f"Table {i+1}")
                self._create_table_slide(prs, table_name, table_data)

        # Metadata slide
        if self.config.include_metadata and data.metadata:
            self._create_metadata_slide(prs, data.metadata)

        # Sources slide
        if self.config.include_sources and data.sources:
            self._create_sources_slide(prs, data.sources)

        # Save to bytes
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()

    def _create_title_slide(self, prs: Any, data: OutputData) -> None:
        """Create title slide."""
        # Use title slide layout (usually index 0)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = data.title

        # Set subtitle if available
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle_text = []

            if data.metadata.get("author"):
                subtitle_text.append(f"By: {data.metadata['author']}")
            if data.metadata.get("date"):
                subtitle_text.append(f"Date: {data.metadata['date']}")

            if subtitle_text:
                subtitle.text = "\n".join(subtitle_text)

    def _create_content_slide(self, prs: Any, title: str, content: str) -> None:
        """Create content slide with title and text."""
        # Use content layout (usually index 1)
        content_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_layout)

        # Set title
        slide.shapes.title.text = title

        # Set content
        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.text = content
            text_frame.word_wrap = True

    def _create_table_slide(self, prs: Any, title: str, table_data: Dict[str, Any]) -> None:
        """Create slide with table."""
        # Use blank or title-only layout
        blank_layout = prs.slide_layouts[5]  # Usually blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Add title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
            top = Inches(1.0)
        else:
            # Add title manually
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9.0)
            height = Inches(0.5)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = title
            p = title_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True
            top = Inches(1.0)

        # Extract table data
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])

        if not headers and not rows:
            return

        # Calculate table dimensions
        num_rows = len(rows) + (1 if headers else 0)
        num_cols = len(headers) if headers else len(rows[0]) if rows else 0

        if num_cols == 0:
            return

        # Add table
        left = Inches(0.5)
        width = Inches(9.0)
        height = Inches(5.0)

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        # Set headers
        if headers:
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                # Make header bold
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True

        # Set data rows
        start_row = 1 if headers else 0
        for row_idx, row in enumerate(rows):
            for col_idx, value in enumerate(row[:num_cols]):
                cell = table.cell(start_row + row_idx, col_idx)
                cell.text = str(value) if value is not None else ""

    def _create_metadata_slide(self, prs: Any, metadata: Dict[str, Any]) -> None:
        """Create metadata slide."""
        content_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_layout)

        slide.shapes.title.text = "Metadata"

        # Build metadata text
        metadata_text = []
        for key, value in metadata.items():
            metadata_text.append(f"{key}: {value}")

        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.text = "\n".join(metadata_text)

    def _create_sources_slide(self, prs: Any, sources: List[Dict[str, Any]]) -> None:
        """Create sources/citations slide."""
        content_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_layout)

        slide.shapes.title.text = "Sources"

        # Build sources text
        sources_text = []
        for i, source in enumerate(sources, 1):
            doc_id = source.get("document_id", "Unknown")
            page = source.get("page_number", "N/A")
            sources_text.append(f"{i}. {doc_id} (Page {page})")

        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.text = "\n".join(sources_text)
